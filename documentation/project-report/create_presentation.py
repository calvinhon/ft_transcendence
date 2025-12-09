#!/usr/bin/env python3
"""
Create comprehensive ft_transcendence presentation based on template structure
Enhanced with detailed technical content and project metrics
"""

from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN
from pptx.dml.color import RGBColor
import os

def load_template():
    """Load the existing template"""
    ppt_path = "Transcendence final.pptx"
    if os.path.exists(ppt_path):
        return Presentation(ppt_path)
    return None

def add_content_to_slide(slide, title, content_bullets):
    """Add title and content to a slide"""
    # Find and update title
    for shape in slide.shapes:
        if hasattr(shape, "text"):
            if shape.is_placeholder:
                phf = shape.placeholder_format
                if phf.type == 1:  # Title
                    if shape.has_text_frame:
                        shape.text_frame.clear()
                        p = shape.text_frame.paragraphs[0]
                        p.text = title
                        p.font.bold = True
                        p.font.size = Pt(44)
            elif "Title and Content" in slide.slide_layout.name:
                if shape.name == "Title 1":
                    shape.text = title

def update_slide_content(prs):
    """Update all slides with enhanced content"""
    
    # Slide 3: Introduction
    slide = prs.slides[2]
    for shape in slide.shapes:
        if hasattr(shape, "text_frame") and "Title" not in shape.name:
            tf = shape.text_frame
            tf.clear()
            
            bullets = [
                "Comprehensive multiplayer Pong game platform",
                "125/125 points compliance with all requirements",
                "Advanced security: JWT, 2FA (TOTP), ModSecurity WAF",
                "Real-time gameplay: 60 FPS server-authoritative sync",
                "Enterprise monitoring: ELK Stack, Prometheus, Grafana",
                "Tournament system with blockchain record-keeping",
                "GDPR compliant user data management",
                "Microservices architecture with Docker deployment"
            ]
            
            for bullet in bullets:
                p = tf.add_paragraph()
                p.text = bullet
                p.level = 0
                p.font.size = Pt(18)
            break
    
    # Slide 4: Major Modules
    slide = prs.slides[3]
    for shape in slide.shapes:
        if hasattr(shape, "text_frame") and "Title" not in shape.name:
            tf = shape.text_frame
            tf.clear()
            
            modules = [
                ("Backend Framework (12/12)", "Fastify 4.29 + Node.js 18+ + TypeScript 5.3"),
                ("Database Connection (12/12)", "SQLite 3 with per-service isolation"),
                ("Backend Gameplay (12/12)", "Server-authoritative Pong with physics"),
                ("Real-Time Sync (12/12)", "WebSocket 60 FPS game state broadcasting"),
                ("OAuth & Authentication (12/12)", "42 School OAuth + JWT + 2FA (TOTP)"),
                ("Blockchain Integration (12/12)", "Hardhat + Solidity for tournament records"),
                ("Server-Side Rendering (12/12)", "SSR service for SEO and offline content"),
            ]
            
            for title, desc in modules:
                p = tf.add_paragraph()
                p.text = f"{title}: {desc}"
                p.level = 0
                p.font.size = Pt(16)
            break
    
    # Slide 5: Minor Modules
    slide = prs.slides[4]
    for shape in slide.shapes:
        if hasattr(shape, "text_frame") and "Title" not in shape.name:
            tf = shape.text_frame
            tf.clear()
            
            modules = [
                "CLI Client (Terminal-based game interface)",
                "AI Opponent (Machine learning-based Pong bot)",
                "Web Application Firewall (ModSecurity + Nginx)",
                "Vault Integration (HashiCorp secrets management)",
                "ELK Logging (Elasticsearch, Logstash, Kibana)",
                "Prometheus Monitoring (Metrics and alerting)",
                "GDPR Compliance (Data export, deletion, consent)",
                "2FA Authentication (TOTP authenticator support)",
                "HTTP-Only Cookies (Secure session management)",
                "Microservices Architecture (4 independent services)",
                "Campaign & Leaderboards (Global rankings system)",
            ]
            
            for module in modules:
                p = tf.add_paragraph()
                p.text = module
                p.level = 0
                p.font.size = Pt(16)
            break
    
    # Slide 6: SDLC
    slide = prs.slides[5]
    for shape in slide.shapes:
        if hasattr(shape, "text_frame") and "Title" not in shape.name:
            tf = shape.text_frame
            tf.clear()
            
            phases = [
                ("Phase 1: Requirements & Planning (Week 1-2)", "Define 125 point requirements, architecture design"),
                ("Phase 2: Core Development (Week 3-7)", "Implement auth, game, user services"),
                ("Phase 3: Advanced Features (Week 8-10)", "Add 2FA, blockchain, monitoring, GDPR"),
                ("Phase 4: Testing & Optimization (Week 11-12)", "180 automated tests, performance tuning"),
                ("Phase 5: Deployment & Documentation (Week 13)", "Docker deployment, final reporting"),
            ]
            
            for phase, desc in phases:
                p = tf.add_paragraph()
                p.text = f"{phase} - {desc}"
                p.level = 0
                p.font.size = Pt(16)
            break
    
    # Slide 8: User Management
    slide = prs.slides[7]
    for shape in slide.shapes:
        if hasattr(shape, "text_frame") and "Title" not in shape.name:
            tf = shape.text_frame
            tf.clear()
            
            content = [
                "Local Registration: Email + secure password (Bcrypt, cost 10)",
                "OAuth 2.0: 42 School intranet integration",
                "Session Management: JWT tokens (HS256, 24h expiry)",
                "User Profiles: Customizable avatars and statistics",
                "Friend System: Add, accept, and manage friend relationships",
                "GDPR Support: Data export and account deletion options",
                "Account Security: Email verification, password reset flow"
            ]
            
            for item in content:
                p = tf.add_paragraph()
                p.text = item
                p.level = 0
                p.font.size = Pt(18)
            break
    
    # Slide 9: Database
    slide = prs.slides[8]
    for shape in slide.shapes:
        if hasattr(shape, "text_frame") and "Title" not in shape.name:
            tf = shape.text_frame
            tf.clear()
            
            content = [
                "Architecture: SQLite per microservice (auth, user, game, tournament)",
                "auth.db: User credentials, sessions, 2FA secrets",
                "users.db: Profiles, friends, statistics, GDPR preferences",
                "games.db: Match history, game events, scores",
                "tournaments.db: Brackets, results, blockchain references",
                "Backup Strategy: Persistent Docker volumes",
                "Encryption: Passwords hashed with Bcrypt, secrets in Vault"
            ]
            
            for item in content:
                p = tf.add_paragraph()
                p.text = item
                p.level = 0
                p.font.size = Pt(18)
            break
    
    # Slide 10: Signing In
    slide = prs.slides[9]
    for shape in slide.shapes:
        if hasattr(shape, "text_frame") and "Title" not in shape.name:
            tf = shape.text_frame
            tf.clear()
            
            items = [
                ("Authentication Methods:", 0),
                ("Local Login: Email + password validation", 1),
                ("OAuth Integration: Seamless 42 School login", 1),
                ("Security Features:", 0),
                ("Password Hashing: Bcrypt with cost 10", 1),
                ("Session Tokens: JWT with 24-hour expiry", 1),
                ("Optional 2FA: TOTP authenticator", 1),
                ("HTTP-Only Cookies: XSS protection", 1),
            ]
            
            for text, level in items:
                p = tf.add_paragraph()
                p.text = text
                p.level = level
                p.font.size = Pt(18)
            break
    
    # Slide 12: Server-Side Rendering
    slide = prs.slides[11]
    for shape in slide.shapes:
        if hasattr(shape, "text_frame") and "Title" not in shape.name:
            tf = shape.text_frame
            tf.clear()
            
            items = [
                ("Purpose: SEO optimization and offline content delivery", 0),
                ("Technology: Node.js with template engine", 0),
                ("Benefits:", 0),
                ("Improved search engine indexing", 1),
                ("Faster initial page load", 1),
                ("Progressive enhancement support", 1),
                ("Reduced client-side processing", 1),
                ("Implementation: Dedicated SSR Service (port 3005)", 0),
            ]
            
            for text, level in items:
                p = tf.add_paragraph()
                p.text = text
                p.level = level
                p.font.size = Pt(18)
            break
    
    # Slide 22: Pong Game
    slide = prs.slides[21]
    for shape in slide.shapes:
        if hasattr(shape, "text_frame") and "Title" not in shape.name:
            tf = shape.text_frame
            tf.clear()
            
            items = [
                ("Game Mechanics:", 0),
                ("2-player Pong: Real-time multiplayer", 1),
                ("Server-Authoritative: No client-side cheating", 1),
                ("60 FPS Synchronization: 16.67ms per frame", 1),
                ("Physics Engine: Ball trajectory and collisions", 1),
                ("Features:", 0),
                ("WebSocket communication for low latency", 1),
                ("Automatic disconnection handling", 1),
                ("Match statistics tracking", 1),
                ("Replay capability (optional)", 1),
            ]
            
            for text, level in items:
                p = tf.add_paragraph()
                p.text = text
                p.level = level
                p.font.size = Pt(17)
            break
    
    # Slide 25: Tournament
    slide = prs.slides[24]
    for shape in slide.shapes:
        if hasattr(shape, "text_frame") and "Title" not in shape.name:
            tf = shape.text_frame
            tf.clear()
            
            items = [
                ("Tournament System:", 0),
                ("Single/Double elimination bracket", 1),
                ("Automatic match scheduling", 1),
                ("Real-time score updates", 1),
                ("Advanced Features:", 0),
                ("Blockchain recording of final results", 1),
                ("Smart contract for tournament integrity", 1),
                ("Leaderboard rankings", 1),
                ("Prize/reward system integration", 1),
            ]
            
            for text, level in items:
                p = tf.add_paragraph()
                p.text = text
                p.level = level
                p.font.size = Pt(18)
            break
    
    # Slide 27: Logging and Monitoring
    slide = prs.slides[26]
    for shape in slide.shapes:
        if hasattr(shape, "text_frame") and "Title" not in shape.name:
            tf = shape.text_frame
            tf.clear()
            
            items = [
                ("Logging Stack (ELK):", 0),
                ("Elasticsearch: Central log storage", 1),
                ("Logstash: Log parsing and enrichment", 1),
                ("Kibana: Dashboard visualization", 1),
                ("Monitoring Stack:", 0),
                ("Prometheus: Metrics collection", 1),
                ("Grafana: Dashboard and alerts", 1),
                ("Custom metrics: API latency, game events, user actions", 1),
            ]
            
            for text, level in items:
                p = tf.add_paragraph()
                p.text = text
                p.level = level
                p.font.size = Pt(18)
            break
    
    # Slide 32: Security
    slide = prs.slides[31]
    for shape in slide.shapes:
        if hasattr(shape, "text_frame") and "Title" not in shape.name:
            tf = shape.text_frame
            tf.clear()
            
            items = [
                ("Network Layer:", 0),
                ("Nginx reverse proxy (ports 80, 443)", 1),
                ("TLS 1.2+ encryption", 1),
                ("ModSecurity WAF protection", 1),
                ("Application Layer:", 0),
                ("JWT authentication (HS256)", 1),
                ("Bcrypt password hashing", 1),
                ("CSRF protection (SameSite cookies)", 1),
                ("Parameterized SQL queries", 1),
                ("XSS prevention (CSP headers)", 1),
                ("Data Protection:", 0),
                ("Vault secrets management", 1),
                ("Database encryption at rest", 1),
                ("HTTP-Only cookies", 1),
            ]
            
            for text, level in items:
                p = tf.add_paragraph()
                p.text = text
                p.level = level
                p.font.size = Pt(16)
            break
    
    # Slide 33: Microservices
    slide = prs.slides[32]
    for shape in slide.shapes:
        if hasattr(shape, "text_frame") and "Title" not in shape.name:
            tf = shape.text_frame
            tf.clear()
            
            items = [
                ("Service Architecture:", 0),
                ("Auth Service (port 3001): User authentication & sessions", 1),
                ("User Service (port 3002): Profiles & social features", 1),
                ("Game Service (port 3003): Pong gameplay & real-time", 1),
                ("Tournament Service (port 3004): Brackets & blockchain", 1),
                ("SSR Service (port 3005): Server-side rendering", 1),
                ("Benefits:", 0),
                ("Independent scaling per service", 1),
                ("Technology flexibility", 1),
                ("Fault isolation", 1),
                ("Continuous deployment capability", 1),
            ]
            
            for text, level in items:
                p = tf.add_paragraph()
                p.text = text
                p.level = level
                p.font.size = Pt(16)
            break
    
    # Slide 34: Limitations
    slide = prs.slides[33]
    for shape in slide.shapes:
        if hasattr(shape, "text_frame") and "Title" not in shape.name:
            tf = shape.text_frame
            tf.clear()
            
            items = [
                ("Single-game: Only Pong implemented", 0),
                ("Max players: 2 per match", 0),
                ("Browser compatibility: Modern browsers only", 0),
                ("Regional deployment: Single server instance", 0),
                ("Blockchain: Ropsten testnet (demo only)", 0),
                ("Storage: SQLite (not production-scale)", 0),
                ("Real-time: ~50ms latency minimum", 0),
            ]
            
            for text, level in items:
                p = tf.add_paragraph()
                p.text = text
                p.level = level
                p.font.size = Pt(18)
            break
    
    # Slide 35: Future Improvements
    slide = prs.slides[34]
    for shape in slide.shapes:
        if hasattr(shape, "text_frame") and "Title" not in shape.name:
            tf = shape.text_frame
            tf.clear()
            
            items = [
                ("Additional games: Chess, Connect 4, others", 0),
                ("Mobile app: Native iOS/Android clients", 0),
                ("Database scaling: PostgreSQL cluster", 0),
                ("Global deployment: Multi-region infrastructure", 0),
                ("Mainnet blockchain: Production smart contracts", 0),
                ("AI improvements: Advanced bot algorithms", 0),
                ("Payment integration: In-game purchases", 0),
                ("Video streaming: Match replay system", 0),
            ]
            
            for text, level in items:
                p = tf.add_paragraph()
                p.text = text
                p.level = level
                p.font.size = Pt(16)
            break

def main():
    """Main execution"""
    prs = load_template()
    
    if prs is None:
        print("❌ Could not load template presentation")
        return
    
    print("📝 Updating presentation with enhanced content...")
    update_slide_content(prs)
    
    output_path = "Transcendence_Enhanced.pptx"
    prs.save(output_path)
    
    print(f"✅ Enhanced presentation created: {output_path}")
    print(f"📊 Total slides: {len(prs.slides)}")
    print(f"💾 File size: {os.path.getsize(output_path) / 1024 / 1024:.1f} MB")

if __name__ == "__main__":
    main()
